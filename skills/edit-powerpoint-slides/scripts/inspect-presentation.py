# /// script
# requires-python = ">=3.12"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Dump the structure of a .pptx so edits can be planned safely.

Reports slide size, the layouts available in the master (needed to add slides
that match the deck), and per slide: layout, placeholders with their idx and
type, shapes, and speaker notes presence. Placeholder idx values are how edits
target the right box, so they are printed explicitly.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


def main() -> int:
    if len(sys.argv) != 2:
        print(f"Usage: uv run {sys.argv[0]} <deck.pptx>")
        return 2
    path = Path(sys.argv[1])
    if not path.exists():
        print(f"Error: file not found: {path}")
        return 1
    if path.suffix.lower() != ".pptx":
        print(f"Error: {path.suffix} is not .pptx; legacy .ppt needs conversion first.")
        return 1

    prs = Presentation(str(path))
    width_in = Emu(prs.slide_width).inches
    height_in = Emu(prs.slide_height).inches
    print(f"Deck: {path}")
    print(f"Slide size: {width_in:.2f} x {height_in:.2f} inches ({'16:9' if abs(width_in / height_in - 16 / 9) < 0.05 else '4:3' if abs(width_in / height_in - 4 / 3) < 0.05 else 'custom'})")
    print(f"Slides: {len(prs.slides)}")

    print("\nLayouts available in master (use these when adding slides):")
    for i, layout in enumerate(prs.slide_layouts):
        ph = [f"idx={p.placeholder_format.idx}:{p.placeholder_format.type}" for p in layout.placeholders]
        print(f"  [{i}] {layout.name}: {ph}")

    for n, slide in enumerate(prs.slides, start=1):
        print(f"\n== Slide {n} (layout: {slide.slide_layout.name!r})")
        for shape in slide.shapes:
            kind = shape.shape_type
            label = f"   shape id={shape.shape_id} name={shape.name!r} type={kind}"
            if shape.is_placeholder:
                label += f" [placeholder idx={shape.placeholder_format.idx} type={shape.placeholder_format.type}]"
            print(label)
            if shape.has_text_frame:
                text = " / ".join(p.text for p in shape.text_frame.paragraphs if p.text)[:90]
                if text:
                    print(f"      text: {text!r}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            print(f"   notes: {slide.notes_slide.notes_text_frame.text.strip()[:90]!r}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
